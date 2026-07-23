import csv
import io
import math
import os
from pathlib import Path


IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.webp'}


def extract_text(file_path: str) -> tuple[str, int]:
    ext = Path(file_path).suffix.lower()
    if ext == '.txt':
        return _extract_txt(file_path)
    elif ext == '.docx':
        return _extract_docx(file_path)
    elif ext == '.pdf':
        return _extract_pdf(file_path)
    elif ext == '.csv':
        return _extract_csv(file_path)
    elif ext == '.xlsx':
        return _extract_xlsx(file_path)
    elif ext == '.pptx':
        return _extract_pptx(file_path)
    elif ext in ('.html', '.htm'):
        return _extract_html(file_path)
    elif ext == '.md':
        return _extract_md(file_path)
    elif ext in IMAGE_EXTENSIONS:
        return _extract_image(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def write_output(text: str, file_path: str, fmt: str, source_path: str = ''):
    os.makedirs(os.path.dirname(file_path), exist_ok=True)
    if fmt == 'txt' or fmt == '':
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(text)
    elif fmt == 'docx':
        from docx import Document
        doc = Document()
        for para in text.split('\n\n'):
            if para.strip():
                doc.add_paragraph(para)
        doc.save(file_path)
    elif fmt == 'csv' and source_path:
        _write_csv(text, file_path, source_path)
    elif fmt == 'xlsx' and source_path:
        _write_xlsx(text, file_path, source_path)
    elif fmt == 'pptx' and source_path:
        _write_pptx(text, file_path, source_path)
    elif fmt in ('html', 'htm') and source_path:
        _write_html(text, file_path, source_path)
    elif fmt == 'md' and source_path:
        _write_md(text, file_path, source_path)
    else:
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(text)


def write_pdf_layout(
    input_path: str,
    output_path: str,
    translate_fn,
    source_lang: str = 'en',
    target_lang: str = 'hi',
) -> dict:
    from .pdf_layout import translate_pdf_preserve
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    return translate_pdf_preserve(
        input_path=input_path,
        output_path=output_path,
        translate_fn=translate_fn,
        source_lang=source_lang,
        target_lang=target_lang,
    )


def write_docx_layout(
    input_path: str,
    output_path: str,
    translate_fn,
    source_lang: str = 'en',
    target_lang: str = 'hi',
) -> dict:
    from .docx_layout import translate_docx_preserve
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    return translate_docx_preserve(
        input_path=input_path,
        output_path=output_path,
        translate_fn=translate_fn,
        source_lang=source_lang,
        target_lang=target_lang,
    )


# ── TXT ──────────────────────────────────────────────────────────────────────

def _extract_txt(file_path: str) -> tuple[str, int]:
    with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
        text = f.read()
    pages = max(1, text.count('\f') + 1)
    return text, pages


# ── DOCX ─────────────────────────────────────────────────────────────────────

def _extract_docx(file_path: str) -> tuple[str, int]:
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = '\n\n'.join(paragraphs)
    return text, len(doc.sections)


# ── PDF ──────────────────────────────────────────────────────────────────────

def _extract_pdf(file_path: str) -> tuple[str, int]:
    import fitz
    doc = fitz.open(file_path)
    total_pages = len(doc)
    text_pages = []
    scanned_pages = []

    for i, page in enumerate(doc):
        page_text = page.get_text().strip()
        if len(page_text) < 30:
            scanned_pages.append(i)
        else:
            text_pages.append((i, page_text))

    has_text = len(text_pages) > total_pages * 0.3
    if has_text and not scanned_pages:
        text = '\n\n'.join(t for _, t in text_pages)
        doc.close()
        return text, total_pages

    ocr_pages = _ocr_pdf_pages(file_path, scanned_pages if scanned_pages else list(range(total_pages)))
    all_text = {}
    for i, t in text_pages:
        all_text[i] = t
    for i, t in ocr_pages.items():
        if i not in all_text:
            all_text[i] = t

    text = '\n\n'.join(all_text.get(i, '') for i in range(total_pages))
    doc.close()
    return text, total_pages


def _ocr_pdf_pages(file_path: str, page_indices: list[int]) -> dict[int, str]:
    import fitz
    import pytesseract
    from PIL import Image
    import io

    doc = fitz.open(file_path)
    ocr_results = {}

    for page_num in page_indices:
        page = doc[page_num]
        pix = page.get_pixmap(dpi=200)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        text = pytesseract.image_to_string(img, lang='hin+guj+eng')
        if text.strip():
            ocr_results[page_num] = text.strip()

    doc.close()
    return ocr_results


# ── CSV ──────────────────────────────────────────────────────────────────────

def _detect_csv_delimiter(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8-sig', errors='replace') as f:
        sample = f.read(4096)
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=',\t;|')
        return dialect.delimiter
    except csv.Error:
        return ','


def _extract_csv(file_path: str) -> tuple[str, int]:
    delimiter = _detect_csv_delimiter(file_path)
    rows = []
    with open(file_path, 'r', encoding='utf-8-sig', errors='replace', newline='') as f:
        reader = csv.reader(f, delimiter=delimiter)
        for row_num, row in enumerate(reader, 1):
            if any(cell.strip() for cell in row):
                cells = [cell.strip() for cell in row]
                rows.append(f"ROW {row_num}: {(' | ').join(cells)}")
    text = '\n'.join(rows)
    pages = max(1, math.ceil(len(rows) / 50))
    return text, pages


def _write_csv(translated_text: str, output_path: str, source_path: str):
    delimiter = _detect_csv_delimiter(source_path)
    parsed_rows = []
    for line in translated_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('ROW '):
            idx = line.index(': ')
            row_data = line[idx + 2:]
            parsed_rows.append(row_data.split(' | '))
        else:
            parsed_rows.append([line])

    with open(output_path, 'w', encoding='utf-8', newline='') as f:
        writer = csv.writer(f, delimiter=delimiter)
        for row in parsed_rows:
            writer.writerow(row)


# ── XLSX ─────────────────────────────────────────────────────────────────────

def _extract_xlsx(file_path: str) -> tuple[str, int]:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True, read_only=True)
    parts = []
    total_rows = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== SHEET: {sheet_name} ===")
        row_num = 0
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_num += 1
                cells = [str(cell).strip() if cell is not None else '' for cell in row]
                parts.append(f"ROW {row_num}: {(' | ').join(cells)}")
        total_rows += row_num

    wb.close()
    text = '\n'.join(parts)
    pages = max(1, math.ceil(total_rows / 30))
    return text, pages


def _write_xlsx(translated_text: str, output_path: str, source_path: str):
    from openpyxl import load_workbook
    wb = load_workbook(source_path)

    sheet_data = {}
    current_sheet = None
    for line in translated_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('=== SHEET: '):
            current_sheet = line[len('=== SHEET: '):-4]
            sheet_data[current_sheet] = []
        elif line.startswith('ROW ') and current_sheet is not None:
            idx = line.index(': ')
            row_data = line[idx + 2:]
            sheet_data[current_sheet].append(row_data.split(' | '))

    for sheet_name, rows in sheet_data.items():
        if sheet_name not in wb.sheetnames:
            continue
        ws = wb[sheet_name]
        for row_idx, row_values in enumerate(rows, 1):
            for col_idx, cell_value in enumerate(row_values, 1):
                if row_idx <= ws.max_row and col_idx <= ws.max_column:
                    ws.cell(row=row_idx, column=col_idx).value = cell_value

    wb.save(output_path)
    wb.close()


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _extract_pptx(file_path: str) -> tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"=== SLIDE: {slide_num} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = '\n'.join(
                    para.text for para in shape.text_frame.paragraphs if para.text.strip()
                )
                if text.strip():
                    parts.append(f"TEXT: {text}")
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows, 1):
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(f"TABLE ROW {row_idx}: {(' | ').join(cells)}")

    text = '\n'.join(parts)
    pages = max(1, len(prs.slides))
    return text, pages


def _write_pptx(translated_text: str, output_path: str, source_path: str):
    from pptx import Presentation
    prs = Presentation(source_path)

    slide_data = {}
    current_slide = None
    for line in translated_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('=== SLIDE: '):
            current_slide = int(line[len('=== SLIDE: '):-3])
            slide_data[current_slide] = {'texts': [], 'tables': []}
        elif line.startswith('TEXT: ') and current_slide is not None:
            slide_data[current_slide]['texts'].append(line[6:])
        elif line.startswith('TABLE ROW ') and current_slide is not None:
            idx = line.index(': ')
            row_data = line[idx + 2:]
            slide_data[current_slide]['tables'].append(row_data.split(' | '))

    for slide_num, data in slide_data.items():
        if slide_num < 1 or slide_num > len(prs.slides):
            continue
        slide = prs.slides[slide_num - 1]
        text_idx = 0
        table_rows = {}
        for shape in slide.shapes:
            if shape.has_text_frame and text_idx < len(data['texts']):
                full_text = data['texts'][text_idx]
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    if para_idx == 0:
                        for run in para.runs:
                            run.text = full_text
                            break
                    else:
                        para.text = ''
                text_idx += 1
            if shape.has_table:
                table_key = id(shape.table)
                table_rows[table_key] = data['tables']

        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                tbl_key = id(tbl)
                if tbl_key in table_rows:
                    for row_idx, row_data in enumerate(table_rows[tbl_key]):
                        if row_idx < len(tbl.rows):
                            for col_idx, cell_val in enumerate(row_data):
                                if col_idx < len(tbl.rows[row_idx].cells):
                                    tbl.rows[row_idx].cells[col_idx].text = cell_val

    prs.save(output_path)


# ── HTML ─────────────────────────────────────────────────────────────────────

_TRANSLATABLE_TAGS = {
    'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
    'p', 'li', 'td', 'th', 'a', 'span',
    'div', 'label', 'button', 'title', 'option',
    'caption', 'summary', 'dt', 'dd', 'blockquote',
}


def _extract_html(file_path: str) -> tuple[str, int]:
    from bs4 import BeautifulSoup
    encodings = ['utf-8', 'latin-1', 'cp1252']
    raw = None
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                raw = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()

    soup = BeautifulSoup(raw, 'html.parser')
    parts = []
    element_count = 0

    for tag in soup.find_all(_TRANSLATABLE_TAGS):
        text = tag.get_text(strip=True)
        if text and not text.startswith('http'):
            tag_name = tag.name
            parts.append(f"<{tag_name}> {text} </{tag_name}>")
            element_count += 1

    text = '\n'.join(parts)
    pages = max(1, math.ceil(element_count / 30))
    return text, pages


def _write_html(translated_text: str, output_path: str, source_path: str):
    from bs4 import BeautifulSoup
    encodings = ['utf-8', 'latin-1', 'cp1252']
    raw = None
    for enc in encodings:
        try:
            with open(source_path, 'r', encoding=enc) as f:
                raw = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        with open(source_path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()

    soup = BeautifulSoup(raw, 'html.parser')
    translated_map = {}
    for line in translated_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        for tag_name in _TRANSLATABLE_TAGS:
            open_tag = f"<{tag_name}> "
            close_tag = f" </{tag_name}>"
            if line.startswith(open_tag) and line.endswith(close_tag):
                inner = line[len(open_tag):-len(close_tag)]
                if tag_name not in translated_map:
                    translated_map[tag_name] = []
                translated_map[tag_name].append(inner)
                break

    for tag_name, translations in translated_map.items():
        idx = 0
        for tag in soup.find_all(tag_name):
            text = tag.get_text(strip=True)
            if text and idx < len(translations):
                tag.clear()
                tag.append(translations[idx])
                idx += 1

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(str(soup))


# ── MARKDOWN ─────────────────────────────────────────────────────────────────

import re

_MD_HEADING_RE = re.compile(r'^(#{1,6})\s+(.+)$')
_MD_LIST_RE = re.compile(r'^(\s*[-*+]|\s*\d+\.)\s+(.+)$')
_MD_TABLE_SEP_RE = re.compile(r'^\|[\s\-:|]+\|$')
_MD_BLOCKQUOTE_RE = re.compile(r'^>\s+(.+)$')
_MD_CODE_FENCE_RE = re.compile(r'^```')
_MD_INLINE_CODE_RE = re.compile(r'`([^`]+)`')
_MD_LINK_RE = re.compile(r'\[([^\]]+)\]\([^\)]+\)')
_MD_IMAGE_RE = re.compile(r'!\[([^\]]*)\]\([^\)]+\)')
_MD_BOLD_RE = re.compile(r'\*\*(.+?)\*\*|__(.+?)__')
_MD_ITALIC_RE = re.compile(r'\*(.+?)\*|_(.+?)_')


def _extract_md(file_path: str) -> tuple[str, int]:
    encodings = ['utf-8', 'latin-1', 'cp1252']
    raw = None
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                raw = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()

    lines = raw.split('\n')
    parts = []
    in_code_block = False
    code_lines = []

    for line in lines:
        if _MD_CODE_FENCE_RE.match(line.strip()):
            if in_code_block:
                parts.append(f"MD_CODE: {' '.join(code_lines)}")
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_lines.append(line)
            continue

        stripped = line.strip()
        if not stripped:
            continue

        heading_match = _MD_HEADING_RE.match(stripped)
        if heading_match:
            level = len(heading_match.group(1))
            text = heading_match.group(2)
            parts.append(f"MD_H{level}: {text}")
            continue

        if _MD_TABLE_SEP_RE.match(stripped):
            continue

        if stripped.startswith('|') and stripped.endswith('|'):
            cells = [c.strip() for c in stripped.split('|')[1:-1]]
            parts.append(f"MD_TABLE: {' | '.join(cells)}")
            continue

        blockquote_match = _MD_BLOCKQUOTE_RE.match(stripped)
        if blockquote_match:
            parts.append(f"MD_QUOTE: {blockquote_match.group(1)}")
            continue

        list_match = _MD_LIST_RE.match(stripped)
        if list_match:
            parts.append(f"MD_LIST: {list_match.group(2)}")
            continue

        clean = stripped
        clean = _MD_IMAGE_RE.sub(r'\1', clean)
        clean = _MD_LINK_RE.sub(r'\1', clean)
        clean = _MD_BOLD_RE.sub(lambda m: m.group(1) or m.group(2), clean)
        clean = _MD_ITALIC_RE.sub(lambda m: m.group(1), clean)
        clean = _MD_INLINE_CODE_RE.sub(r'\1', clean)

        if clean.strip():
            parts.append(f"MD_PARA: {clean.strip()}")

    text = '\n'.join(parts)
    pages = max(1, math.ceil(len(parts) / 40))
    return text, pages


def _write_md(translated_text: str, output_path: str, source_path: str):
    output_lines = []
    for line in translated_text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('MD_H1: '):
            output_lines.append(f'# {line[7:]}')
        elif line.startswith('MD_H2: '):
            output_lines.append(f'## {line[7:]}')
        elif line.startswith('MD_H3: '):
            output_lines.append(f'### {line[7:]}')
        elif line.startswith('MD_H4: '):
            output_lines.append(f'#### {line[7:]}')
        elif line.startswith('MD_H5: '):
            output_lines.append(f'##### {line[7:]}')
        elif line.startswith('MD_H6: '):
            output_lines.append(f'###### {line[7:]}')
        elif line.startswith('MD_LIST: '):
            output_lines.append(f'- {line[9:]}')
        elif line.startswith('MD_QUOTE: '):
            output_lines.append(f'> {line[10:]}')
        elif line.startswith('MD_TABLE: '):
            cells = line[10:].split(' | ')
            output_lines.append('| ' + ' | '.join(cells) + ' |')
        elif line.startswith('MD_CODE: '):
            output_lines.append(f'```\n{line[9:]}\n```')
        elif line.startswith('MD_PARA: '):
            output_lines.append(line[9:])
        else:
            output_lines.append(line)

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(output_lines) + '\n')


# ── IMAGES ───────────────────────────────────────────────────────────────────

def _extract_image(file_path: str) -> tuple[str, int]:
    import pytesseract
    from PIL import Image

    img = Image.open(file_path)

    if img.mode != 'RGB':
        img = img.convert('RGB')

    max_dim = 3000
    if max(img.size) > max_dim:
        ratio = max_dim / max(img.size)
        new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    text = pytesseract.image_to_string(img, lang='hin+guj+eng')
    img.close()

    if not text.strip():
        return "[No text detected in image]", 1

    return text.strip(), 1
