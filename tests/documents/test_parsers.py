import os
import csv
import sys
import tempfile
import pytest
from unittest.mock import patch, MagicMock

docx = pytest.importorskip("docx")
openpyxl = pytest.importorskip("openpyxl")
pptx = pytest.importorskip("pptx")

from apps.documents.parsers import (
    extract_text,
    _extract_txt,
    _extract_docx,
    _extract_pdf,
    _extract_csv,
    _extract_xlsx,
    _extract_pptx,
    _extract_html,
    _write_csv,
    _write_xlsx,
    _write_pptx,
    _write_html,
    write_output,
    _detect_csv_delimiter,
)

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation


@pytest.fixture
def txt_file():
    path = os.path.join(tempfile.gettempdir(), 'parser_test.txt')
    with open(path, 'w') as f:
        f.write("Hello world\nSecond line\n\nThird paragraph")
    return path


@pytest.fixture
def csv_file():
    path = os.path.join(tempfile.gettempdir(), 'parser_test.csv')
    with open(path, 'w', newline='') as f:
        writer = csv.writer(f)
        writer.writerow(['Name', 'Age', 'City'])
        writer.writerow(['Alice', '30', 'Delhi'])
        writer.writerow(['Bob', '25', 'Mumbai'])
        writer.writerow(['Charlie', '35', 'Bangalore'])
    return path


@pytest.fixture
def csv_file_tab():
    path = os.path.join(tempfile.gettempdir(), 'parser_test_tab.csv')
    with open(path, 'w') as f:
        f.write("Name\tAge\tCity\n")
        f.write("Alice\t30\tDelhi\n")
        f.write("Bob\t25\tMumbai\n")
    return path


@pytest.fixture
def html_file():
    path = os.path.join(tempfile.gettempdir(), 'parser_test.html')
    content = """<!DOCTYPE html>
<html>
<head><title>Test Page</title></head>
<body>
    <h1>Welcome</h1>
    <p>This is a paragraph.</p>
    <ul>
        <li>Item one</li>
        <li>Item two</li>
    </ul>
    <a href="https://example.com">Click here</a>
</body>
</html>"""
    with open(path, 'w') as f:
        f.write(content)
    return path


@pytest.fixture
def xlsx_file():
    path = os.path.join(tempfile.gettempdir(), 'parser_test.xlsx')
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(['Name', 'Age', 'City'])
    ws.append(['Alice', 30, 'Delhi'])
    ws.append(['Bob', 25, 'Mumbai'])
    wb.save(path)
    wb.close()
    return path


@pytest.fixture
def pptx_file():
    path = os.path.join(tempfile.gettempdir(), 'parser_test.pptx')
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Hello World"
    subtitle = slide.placeholders[1]
    subtitle.text = "Subtitle text here"
    prs.save(path)
    return path


# ── TXT Parser ────────────────────────────────────────────────────────────────


class TestExtractTxt:
    def test_basic_text(self, txt_file):
        text, pages = _extract_txt(txt_file)
        assert "Hello world" in text
        assert "Second line" in text
        assert pages == 1

    def test_page_count_with_formfeed(self):
        path = os.path.join(tempfile.gettempdir(), 'parser_ff.txt')
        with open(path, 'w') as f:
            f.write("Page one\n\nPage two\n\nPage three")
        text, pages = _extract_txt(path)
        assert pages == 1

    def test_empty_file(self):
        path = os.path.join(tempfile.gettempdir(), 'empty.txt')
        with open(path, 'w') as f:
            f.write("")
        text, pages = _extract_txt(path)
        assert text == ""
        assert pages == 1

    def test_dispatch(self, txt_file):
        text, pages = extract_text(txt_file)
        assert "Hello world" in text
        assert pages >= 1


# ── DOCX Parser ───────────────────────────────────────────────────────────────


class TestExtractDocx:
    def test_docx_basic(self):
        path = os.path.join(tempfile.gettempdir(), 'parser_test.docx')
        doc = Document()
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")
        doc.add_paragraph("Third paragraph")
        doc.save(path)

        text, pages = _extract_docx(path)
        assert "First paragraph" in text
        assert "Second paragraph" in text
        assert "Third paragraph" in text
        assert pages >= 1

    def test_docx_empty_paragraphs_filtered(self):
        path = os.path.join(tempfile.gettempdir(), 'parser_empty_para.docx')
        doc = Document()
        doc.add_paragraph("")
        doc.add_paragraph("Real text")
        doc.add_paragraph("")
        doc.save(path)

        text, pages = _extract_docx(path)
        assert "Real text" in text
        assert text.strip().count('\n\n') == 0

    def test_docx_dispatch(self):
        path = os.path.join(tempfile.gettempdir(), 'dispatch_test.docx')
        doc = Document()
        doc.add_paragraph("Dispatch test")
        doc.save(path)

        text, pages = extract_text(path)
        assert "Dispatch test" in text
        assert pages >= 1


# ── PDF Parser ────────────────────────────────────────────────────────────────


class TestExtractPdf:
    def _make_mock_fitz(self, pages):
        mock_fitz = MagicMock()
        mock_pages = []
        for text in pages:
            p = MagicMock()
            p.get_text.return_value = text
            mock_pages.append(p)
        mock_doc = MagicMock()
        mock_doc.__iter__ = lambda s: iter(mock_pages)
        mock_doc.__len__ = lambda s: len(mock_pages)
        mock_fitz.open.return_value = mock_doc
        return mock_fitz

    def test_pdf_text_dominant(self):
        mock_fitz = self._make_mock_fitz(["A" * 100, "B" * 100])
        with patch.dict(sys.modules, {'fitz': mock_fitz}):
            text, pages = _extract_pdf('dummy.pdf')
        assert pages == 2
        assert "A" * 100 in text
        assert "B" * 100 in text

    @patch('apps.documents.parsers._ocr_pdf_pages')
    def test_pdf_scanned_pages(self, mock_ocr):
        mock_fitz = self._make_mock_fitz(["short"])
        mock_ocr.return_value = {0: "OCR extracted text"}
        with patch.dict(sys.modules, {'fitz': mock_fitz}):
            text, pages = _extract_pdf('dummy.pdf')
        assert pages == 1
        assert "OCR extracted text" in text

    @patch('apps.documents.parsers._ocr_pdf_pages')
    def test_pdf_mixed_text_and_scanned(self, mock_ocr):
        mock_fitz = self._make_mock_fitz(["A" * 100, "short"])
        mock_ocr.return_value = {1: "Scanned content"}
        with patch.dict(sys.modules, {'fitz': mock_fitz}):
            text, pages = _extract_pdf('dummy.pdf')
        assert pages == 2
        assert "A" * 100 in text
        assert "Scanned content" in text


# ── CSV Parser ────────────────────────────────────────────────────────────────


class TestExtractCsv:
    def test_csv_basic(self, csv_file):
        text, pages = _extract_csv(csv_file)
        assert "ROW 1:" in text
        assert "Name" in text
        assert "Alice" in text
        assert "Delhi" in text
        assert pages >= 1

    def test_csv_tab_delimited(self, csv_file_tab):
        text, pages = _extract_csv(csv_file_tab)
        assert "ROW 1:" in text
        assert "Alice" in text

    def test_csv_page_count_scales(self):
        path = os.path.join(tempfile.gettempdir(), 'big.csv')
        with open(path, 'w', newline='') as f:
            writer = csv.writer(f)
            for i in range(100):
                writer.writerow([f'Row {i}', f'Value {i}'])
        text, pages = _extract_csv(path)
        assert pages > 1

    def test_csv_empty_rows_filtered(self):
        path = os.path.join(tempfile.gettempdir(), 'empty_rows.csv')
        with open(path, 'w', newline='') as f:
            writer = csv.writer(f)
            writer.writerow(['A', 'B'])
            writer.writerow(['', ''])
            writer.writerow(['C', 'D'])
        text, pages = _extract_csv(path)
        assert "A" in text
        assert "C" in text
        lines = [l for l in text.split('\n') if l.strip()]
        assert len(lines) == 2

    def test_csv_dispatch(self, csv_file):
        text, pages = extract_text(csv_file)
        assert "ROW 1:" in text
        assert pages >= 1

    def test_detect_csv_delimiter_comma(self):
        path = os.path.join(tempfile.gettempdir(), 'comma.csv')
        with open(path, 'w') as f:
            f.write("a,b,c\n1,2,3\n")
        assert _detect_csv_delimiter(path) == ','

    def test_detect_csv_delimiter_tab(self):
        path = os.path.join(tempfile.gettempdir(), 'tab.csv')
        with open(path, 'w') as f:
            f.write("a\tb\tc\n1\t2\t3\n")
        assert _detect_csv_delimiter(path) == '\t'


# ── XLSX Parser ───────────────────────────────────────────────────────────────


class TestExtractXlsx:
    def test_xlsx_basic(self, xlsx_file):
        text, pages = _extract_xlsx(xlsx_file)
        assert "=== SHEET: Sheet1 ===" in text
        assert "Name" in text
        assert "Alice" in text
        assert pages >= 1

    def test_xlsx_multiple_sheets(self):
        path = os.path.join(tempfile.gettempdir(), 'multi.xlsx')
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Data"
        ws1.append(['Col1', 'Col2'])
        ws1.append(['val1', 'val2'])
        ws2 = wb.create_sheet("Summary")
        ws2.append(['Total'])
        ws2.append(['100'])
        wb.save(path)
        wb.close()

        text, pages = _extract_xlsx(path)
        assert "=== SHEET: Data ===" in text
        assert "=== SHEET: Summary ===" in text
        assert "val1" in text
        assert "100" in text

    def test_xlsx_dispatch(self, xlsx_file):
        text, pages = extract_text(xlsx_file)
        assert "=== SHEET:" in text
        assert pages >= 1


# ── PPTX Parser ───────────────────────────────────────────────────────────────


class TestExtractPptx:
    def test_pptx_basic(self, pptx_file):
        text, pages = _extract_pptx(pptx_file)
        assert "=== SLIDE: 1 ===" in text
        assert "Hello World" in text
        assert pages >= 1

    def test_pptx_empty_slide(self):
        path = os.path.join(tempfile.gettempdir(), 'empty_slide.pptx')
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        prs.slides.add_slide(slide_layout)
        prs.save(path)

        text, pages = _extract_pptx(path)
        assert pages == 1
        assert "=== SLIDE: 1 ===" in text

    def test_pptx_dispatch(self, pptx_file):
        text, pages = extract_text(pptx_file)
        assert "=== SLIDE:" in text
        assert pages >= 1


# ── HTML Parser ───────────────────────────────────────────────────────────────


class TestExtractHtml:
    def test_html_basic(self, html_file):
        text, pages = _extract_html(html_file)
        assert "<h1> Welcome </h1>" in text
        assert "<p> This is a paragraph. </p>" in text
        assert "<li> Item one </li>" in text
        assert pages >= 1

    def test_html_skips_urls(self):
        path = os.path.join(tempfile.gettempdir(), 'url.html')
        with open(path, 'w') as f:
            f.write("<html><body><a href='https://example.com'>https://example.com</a></body></html>")
        text, pages = _extract_html(path)
        assert "https://example.com" not in text

    def test_html_page_count_scales(self):
        path = os.path.join(tempfile.gettempdir(), 'big.html')
        items = ''.join(f'<p>Item {i}</p>' for i in range(100))
        with open(path, 'w') as f:
            f.write(f"<html><body>{items}</body></html>")
        text, pages = _extract_html(path)
        assert pages > 1

    def test_html_dispatch(self, html_file):
        text, pages = extract_text(html_file)
        assert "<h1>" in text
        assert pages >= 1

    def test_html_htm_extension(self):
        path = os.path.join(tempfile.gettempdir(), 'test.htm')
        with open(path, 'w') as f:
            f.write("<html><body><p>HTM content</p></body></html>")
        text, pages = extract_text(path)
        assert "HTM content" in text


# ── Write Output ──────────────────────────────────────────────────────────────


class TestWriteOutput:
    def test_write_txt(self):
        out = os.path.join(tempfile.gettempdir(), 'write_out.txt')
        write_output("Hello world\nSecond line", out, 'txt')
        with open(out) as f:
            assert f.read() == "Hello world\nSecond line"

    def test_write_txt_empty_format(self):
        out = os.path.join(tempfile.gettempdir(), 'write_out_empty.txt')
        write_output("Fallback text", out, '')
        with open(out) as f:
            assert f.read() == "Fallback text"

    def test_write_docx(self):
        out = os.path.join(tempfile.gettempdir(), 'write_out.docx')
        write_output("First para\n\nSecond para", out, 'docx')
        doc = Document(out)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        assert "First para" in paragraphs
        assert "Second para" in paragraphs

    def test_write_csv(self):
        src = os.path.join(tempfile.gettempdir(), 'src.csv')
        with open(src, 'w', newline='') as f:
            writer = csv.writer(f)
            writer.writerow(['A', 'B'])
            writer.writerow(['1', '2'])

        out = os.path.join(tempfile.gettempdir(), 'write_out.csv')
        translated = "ROW 1: X | Y"
        _write_csv(translated, out, src)

        with open(out) as f:
            reader = csv.reader(f)
            rows = list(reader)
        assert rows[0] == ['X', 'Y']

    def test_write_pptx(self):
        src = os.path.join(tempfile.gettempdir(), 'src.pptx')
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Original"
        prs.save(src)

        out = os.path.join(tempfile.gettempdir(), 'write_out.pptx')
        translated = "=== SLIDE: 1 ===\nTEXT: Translated title"
        _write_pptx(translated, out, src)

        prs2 = Presentation(out)
        assert prs2.slides[0].shapes.title.text == "Translated title"

    def test_write_html(self):
        src = os.path.join(tempfile.gettempdir(), 'src.html')
        with open(src, 'w') as f:
            f.write("<html><body><h1>Hello</h1><p>World</p></body></html>")

        out = os.path.join(tempfile.gettempdir(), 'write_out.html')
        translated = "<h1> translated </h1>\n<p> translated p </p>"
        _write_html(translated, out, src)

        with open(out) as f:
            content = f.read()
        assert "translated" in content

    def test_write_xlsx(self):
        src = os.path.join(tempfile.gettempdir(), 'src.xlsx')
        wb = Workbook()
        ws = wb.active
        ws.append(['Name', 'Value'])
        ws.append(['Alice', '100'])
        wb.save(src)
        wb.close()

        out = os.path.join(tempfile.gettempdir(), 'write_out.xlsx')
        translated = "=== SHEET: Sheet ===\nROW 1: translated_name | value_translated"
        _write_xlsx(translated, out, src)

        wb2 = load_workbook(out)
        ws2 = wb2.active
        assert ws2.cell(1, 1).value == 'translated_name'
        assert ws2.cell(1, 2).value == 'value_translated'
        wb2.close()

    def test_write_unsupported_format_falls_back(self):
        out = os.path.join(tempfile.gettempdir(), 'write_out.xyz')
        write_output("Fallback", out, 'xyz')
        with open(out) as f:
            assert f.read() == "Fallback"


# ── Unsupported File ──────────────────────────────────────────────────────────


class TestUnsupportedFormat:
    def test_unsupported_extension_raises(self):
        path = os.path.join(tempfile.gettempdir(), 'test.xyz')
        with open(path, 'w') as f:
            f.write("data")
        with pytest.raises(ValueError, match="Unsupported file format"):
            extract_text(path)
